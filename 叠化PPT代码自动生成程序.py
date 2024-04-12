from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide_with_text_at_bottom_left(presentation, text):
    slide_layout = presentation.slide_layouts[5]  
    slide = presentation.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    left = Inches(-1.1)  # 文本框距离PPT左侧的距离，该数字需视情况调整
    top = Inches(6)      # 文本框距离PPT上侧的距离，该数字需视情况调整
    width = Inches(9)    # 文本框宽度，大概率不需要调整
    height = Inches(1)   # 文本框长度，大概率不需要调整
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(36) #设置字体大小为36，视情况修改

    p.font.color.rgb = RGBColor(255, 255, 255) 
    p.alignment = PP_ALIGN.LEFT

    return presentation

# 该PPT的最终存储位置。注意1.一定要用双\\，2.末尾一定是该PPT的名字，且加后缀.pptx
output_file_path = 'C:\\Users\\89377\\Desktop\\cjcsb.pptx'

# 歌词文档，每一行生成一页PPT，空格行生成一页空白页
lyrics = """
千山万水相聚的一瞬
千言万语就在一个眼神
生活是个复杂的剧本
不改变我们生命的单纯

不问扬起过多少烟尘
不枉内心一直追求的安顿
不管走过多远的旅程

感动不一定流泪
感情还一样率真

我为你留着一盏灯
让你心境永远不会近黄昏

我心中不会有黄昏
有你在永远像初春的清晨

云很淡
云很淡
风很清
风很清
任星辰
任星辰
浮浮沉沉
生活是个复杂的剧本
不改变我们生命的单纯
不管走过多远的旅程
感动不一定流泪
感情还一样率真
我心中亮着一盏灯
你是让我看透天地那个人
你是我心里那盏灯
让我静看外面喧闹的红尘
且听岁月像旋律永恒
一直陪伴不断聚散的旅程
我心中开着一扇门
一直等待永远青春的归人
云很淡
云很淡
风很清
风很清
任星辰
浮浮沉沉
"""

lines = lyrics.strip().split('\n')
presentation = Presentation()
for line in lines:
    presentation = add_slide_with_text_at_bottom_left(presentation, line)
presentation.save(output_file_path)